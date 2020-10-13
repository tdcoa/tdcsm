"Powerpoint instantiation from a template"
from __future__ import annotations

import csv
import logging
import re
from argparse import ArgumentParser
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import (Any, ClassVar, Dict, Iterable, List, Optional, Tuple, Type,
					Union)

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import GroupShapes, SlideShapes
from pptx.text.text import TextFrame, Font

ShapeContainer = Union[SlideShapes, GroupShapes]
logging.basicConfig(format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)
tags: Dict[str, Type[Placeholder]] = {}


@dataclass
class Location:
	"Location of pattern"
	slnum: int                     # slide number
	container: ShapeContainer      # parent container of the shape
	shape: BaseShape               # shape
	tbl_row: Optional[int] = None  # for table shapes, the row
	tbl_col: Optional[int] = None  # for table shapes, the col

	def __str__(self) -> str:
		return f"slide#={self.slnum}, shape={self.shape.name}, row={self.tbl_row}, col={self.tbl_col}"

	@property
	def ftext(self) -> TextFrame:
		"formatted text"
		return self.shape.text_frame if self.tbl_row is None else self.shape.table.cell(self.tbl_row, self.tbl_col).text_frame

	def ftext_iter(self) -> Iterable[TextFrame]:
		"iterate over cells formatted text attributes in the column"
		yield from (r.cells[self.tbl_col].text_frame for r in self.shape.table.rows)


@dataclass
class Placeholder:
	"Base, Abastract class for place-holders"
	tag: ClassVar[str] = '?'

	loc: Location
	datafile: str
	pat: str

	def __str__(self) -> str:
		return f"type={self.tag}, data={self.datafile}, pat={self.pat}, location=({self.loc})"

	def replace(self, datapath: Path) -> None:
		"replace placeholders with values loaded from the datapath"
		raise NotImplementedError("cannot call abstract method")

	@classmethod
	def parse(cls, loc: Location, data: str, pat: str) -> Placeholder:
		"parse values to instantiate an object"
		return cls(loc, data, pat)

	@classmethod
	def __init_subclass__(cls, **_: Any) -> None:
		tags[cls.tag] = cls


@dataclass
class PicPlaceHolder(Placeholder):
	"Placeholder for a Picture"
	tag: ClassVar[str] = 'pic'

	def replace(self, datapath: Path) -> None:
		logger.debug("Replacing: %s", self.datafile)
		self.loc.container.add_picture(
			str(datapath / self.datafile),
			left=self.loc.shape.left,
			top=self.loc.shape.top,
			height=self.loc.shape.height,
			width=self.loc.shape.width
		)
		self.loc.shape.text = f'**{datapath / self.datafile}**'

	@classmethod
	def parse(cls, loc: Location, data: str, pat: str) -> Placeholder:
		return cls(loc, data, pat)


@dataclass
class ColPlaceHolder(Placeholder):
	"Placeholder for a column of values from a csv file"
	tag: ClassVar[str] = 'col'

	colnum: int

	def replace(self, datapath: Path) -> None:
		logger.debug("Replacing: %s[%d]", self.datafile, self.colnum)
		data = (r[self.colnum - 1] for r in load_csv(datapath / self.datafile))
		font = None
		for tf, datum in zip(self.loc.ftext_iter(), data):
			font = repl_text(tf, None, str(datum), font)

	@classmethod
	def parse(cls, loc: Location, data: str, pat: str) -> Placeholder:
		m = re.fullmatch(r"([^\[\]]+)\[(\d+)\]", data)
		if m is None:
			raise ValueError("valid 'col' specification must '<data>[<colnum>]'")
		return cls(loc, m.group(1), pat, int(m.group(2)))


@dataclass
class ValPlaceHolder(Placeholder):
	"Placeholder for a single value from a csv file"
	tag: ClassVar[str] = 'val'

	rownum: int
	colnum: int

	def replace(self, datapath: Path) -> None:
		try:
			data = load_csv(datapath / self.datafile)[self.rownum][self.colnum - 1]
		except IndexError:
			data = None

		if data is None:
			raise ValueError(f'Invalid cell#[{self.rownum}:{self.colnum}] in {datapath / self.datafile}')

		repl_text(self.loc.ftext, self.pat, data)

	@classmethod
	def parse(cls, loc: Location, data: str, pat: str) -> Placeholder:
		m = re.fullmatch(r"([^\[\]]+)\[(\d+):(\d+)\]", data)
		if m is None:
			raise ValueError(f"{data} is invalid 'val' specification; must '<data>[<rownum>:<colnum>]'")
		return cls(loc, m.group(1), pat, int(m.group(2)), int(m.group(3)))


def repl_text(tf: TextFrame, src: Optional[str], tgt: str, dflt_font: Optional[Font] = None) -> Optional[Font]:
	"replace src text with tgt text in a shape's text_frame object"
	logger.debug("repl_text() %s with %s", src, tgt)
	fnm = fsz = fbd = fil = fco = None

	def save_style(font: Font) -> None:
		"save current font properties"
		nonlocal fnm, fsz, fbd, fil, fco

		fnm = font.name or fnm
		fsz = font.size or fsz
		fbd = font.bold or fbd
		fil = font.italic or fil
		fco = font.color or fco
		logger.debug("repl_text(), saving style: Font Name: %s, Size: %s, Bold: %s, Ital: %s", fnm, fsz, fbd, fil)

	def reapply_style(font: Font) -> None:
		"reapply saved font properties"
		logger.debug("repl_text(), repplying style: Font Name: %s, Size: %s, Bold: %s, Ital: %s", fnm, fsz, fbd, fil)
		if fnm is not None:
			font.name = fnm
		if fsz is not None:
			font.size = fsz
		if fbd is not None:
			font.bold = fbd
		if fil is not None:
			font.italic = fil
		if fco is not None:
			if fco.type == 1:
				font.color.rgb = fco.rgb
			elif fco.type == 2:
				font.color.theme_color = fco.theme_color
			elif fco.type is not None:
				logger.debug('Unknown color type: %d', fco.type)

	if dflt_font is not None:
		save_style(dflt_font)

	for para in tf.paragraphs:
		if src is None or src in para.text:  # src == None => replace any value
			save_style(para.font)
			for e2, r in enumerate(para.runs):
				logger.debug("repl_text(), found run# %d", e2)
				save_style(r.font)
				break
			if src is None:
				para.text = tgt
			else:
				para.text = para.text.replace(src, tgt)
			reapply_style(para.font)
			return para.font

	return None


@lru_cache
def load_csv(csvpath: Path) -> List[List[Any]]:
	"load CSV file and return a list of all rows"
	with csvpath.open() as f:
		return list(csv.reader(f))


def iter_shapes(container: ShapeContainer) -> Iterable[Tuple[ShapeContainer, BaseShape]]:
	"iterate recursively over member shapes returning shapes's immediate container and the shape"
	for sh in container:
		if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
			yield from iter_shapes(sh.shapes)
		else:
			yield (container, sh)


def findtag_iter(shape: BaseShape) -> Iterable[Tuple[re.Match, Optional[int], Optional[int]]]:
	"find and iterate over all match objects matching tags in a shape"
	if shape.has_table:
		for rnum, r in enumerate(shape.table.rows):
			for cnum, c in enumerate(r.cells):
				m2 = re.fullmatch("{{([^:]+):(.+?)}}", c.text.rstrip())
				if m2:
					yield (m2, rnum, cnum)

	elif shape.has_text_frame:
		for m in re.finditer("{{([^:]+):(.+?)}}", shape.text):
			yield (m, None, None)


def find_placeholders(ppt: Presentation) -> Iterable[Placeholder]:
	"iterate (shape, palceholder) pair over all presentation placeholders"
	for slnum, slide in enumerate(ppt.slides, start=1):
		for container, shape in iter_shapes(slide.shapes):
			logger.debug(f"slide#: {slnum}, type: {shape.shape_type}, text?: {shape.has_text_frame}, table?: {shape.has_table}")
			for m, rnum, cnum in findtag_iter(shape):
				tag, datafile, pat = m.group(1), m.group(2), m.group(0)
				try:
					yield tags[tag].parse(Location(slnum, container, shape, rnum, cnum), datafile, pat)
				except KeyError as ex:
					logger.error("Invalid replacement tag '%s' on slide# %d, shape name: %s", str(ex), slnum, shape.name)
				except ValueError as ex:
					logger.error("Invalid Tag[slide#=%d, shape name=%s]: %s", slnum, shape.name, str(ex))


def replace_placeholders(pptpath: Path, datapath: Path, output: Optional[Path] = None) -> None:
	"replace all place-holders with the contents from the file"
	load_csv.cache_clear()  # cache only for this call

	logger.debug("Starting replacing placeholders in %s", pptpath)
	ppt = pptx.Presentation(pptpath)
	for ph in find_placeholders(ppt):
		try:
			ph.replace(datapath)
		except FileNotFoundError as err:
			logger.warning("Ignoring missing '%s', while creating '%s'", err.filename, pptpath.name)
	ppt.save(output if output is not None else pptpath)


def run(ppttmpl: Path, pptout: Optional[Path] = None, data: Optional[Path] = None, debug: bool = False) -> None:
	"script entry-point"
	if debug:
		logger.setLevel(logging.DEBUG)

	if pptout is None:
		ppt = pptx.Presentation(ppttmpl)
		for ph in find_placeholders(ppt):
			print(str(ph))
	else:
		replace_placeholders(ppttmpl, (data if data is not None else ppttmpl.parent), output=pptout)


if __name__ == '__main__':
	p = ArgumentParser(description=__doc__)

	p.add_argument("ppttmpl", type=Path, help="Presentation template that contains place-holder markers")
	p.add_argument("pptout", type=Path, nargs='?', help="Perform replacements and write output")
	p.add_argument("--data", type=Path, help="folder containing substitution data, defaults to same path as template")
	p.add_argument("--debug", action='store_true', help="show debug messages")

	run(**vars(p.parse_args()))
