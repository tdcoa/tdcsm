"test cases for powerpoint templating"
from typing import List, Any
from pathlib import Path
import pytest
import pptx
from tdcsm.pptx import find_placeholders, replace_placeholders


@pytest.fixture(scope="session")
def pptxdir(testdir: Path, tmp_path_factory: Any) -> Path:
	"return a temporary directory containing copies of data files"
	tmpdir = tmp_path_factory.mktemp("tdcsm_pptx")
	for p in (testdir / "test_pptx").iterdir():
		(tmpdir / p.name).write_bytes(p.read_bytes())
	return tmpdir


def phlist(fname: Path) -> List[str]:
	"return string representations of all place-holders in the specified ppt file"
	return [str(ph) for ph in find_placeholders(pptx.Presentation(fname))]


def test_ppt_pic(pptxdir: Path) -> None:
	"assert pic type placeholder exists"
	assert phlist(pptxdir / "test_pic.pptx") == [
		'type=pic, data=alldates.png, pat={{pic:alldates.png}}, location=(slide#=1, shape=Rectangle 3, row=None, col=None)'
	]


def test_ppt_col(pptxdir: Path) -> None:
	"assert table column type placeholders exist"
	assert phlist(pptxdir / "test_col.pptx") == [
		'type=col, data=dates.csv, pat={{col:dates.csv[1]}}, location=(slide#=1, shape=Table 4, row=0, col=0)',
		'type=col, data=dates.csv, pat={{col:dates.csv[2]}}, location=(slide#=1, shape=Table 4, row=0, col=1)'
	]


def test_ppt_val(pptxdir: Path) -> None:
	"assert val type placeholder exists"
	assert phlist(pptxdir / "test_val.pptx") == [
		'type=val, data=birthday.csv, pat={{val:birthday.csv[1:2]}}, location=(slide#=1, shape=Content Placeholder 1, row=None, col=None)'
	]


def test_ppt_multival(pptxdir: Path) -> None:
	"assert multiple paragraphss have same the font"
	replace_placeholders(pptxdir / "test_multiline.pptx", pptxdir)
	ppt = pptx.Presentation(pptxdir / "test_multiline.pptx")
	shapes = [sh for sn, sl in enumerate(ppt.slides, start=1) for sh in sl.shapes]

	assert [para.font.name for para in shapes[0].shapes[0].text_frame.paragraphs] == ['Arial', 'Arial', 'Arial']


def test_missing_pic(pptxdir: Path) -> None:
	"assert missing pic file doesn't cause an exception"
	replace_placeholders(pptxdir / "test_missing_pic.pptx", pptxdir)


def test_missing_csv(pptxdir: Path) -> None:
	"assert missing CSV file doesn't cause an exception"
	replace_placeholders(pptxdir / "test_missing_csv.pptx", pptxdir)


def test_ppt_all(pptxdir: Path) -> None:
	"assert values being replaced are same as source"
	replace_placeholders(pptxdir / "test_all.pptx", pptxdir)
	ppt = pptx.Presentation(pptxdir / "test_all.pptx")
	shapes = [sh for sn, sl in enumerate(ppt.slides, start=1) for sh in sl.shapes]

	assert shapes[3].image.blob == (pptxdir / "alldates.png").read_bytes()
	assert shapes[4].text == 'Text with one placeholder 1974-12-17'
	assert [[shapes[6].table.cell(r, c).text for c in range(2)] for r in range(3)] == [
		['cal_date', 'item'],
		['2020/01/01', 'Decade'],
		['2000/01/01', 'Millennium']
	]
	assert shapes[7].text == 'TextBox with multiple placehoders: USA received independence on 1776/07/04'
